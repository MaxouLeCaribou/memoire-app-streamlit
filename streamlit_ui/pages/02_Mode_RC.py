# streamlit_ui/pages/02_Mode_RC.py
import io, os, sys
import streamlit as st
import re
import uuid
import datetime
import json
from io import BytesIO
import fitz  # PyMuPDF
import time
import traceback
import chromadb
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np



# --- FOND ÉCRAN via base64 (fiable) ---
import base64, mimetypes
from pathlib import Path

def set_background(img_path: str):
    p = Path(img_path)
    if not p.is_file():
        import streamlit as st
        st.warning(f"Image introuvable : {p.resolve()}")
        return
    mime = mimetypes.guess_type(p.name)[0] or "image/jpg"
    data = base64.b64encode(p.read_bytes()).decode()

    import streamlit as st
    st.markdown(f"""
    <style>
    [data-testid="stAppViewContainer"] {{
      background-image: url("data:{mime};base64,{data}");
      background-repeat: no-repeat;
      background-position: center center;
      background-attachment: fixed;
      background-size: cover;
      position: relative;
    }}

    [data-testid="stAppViewContainer"]::before {{
      content: "";
      position: absolute;
      top: 0; left: 0; right: 0; bottom: 0;
      background: rgba(0,0,0,0.8); /* assombrissement léger */
      z-index: 0;
    }}

    /* contenu au-dessus du voile */
    [data-testid="stAppViewContainer"] * {{
      position: relative;
      z-index: 1;
    }}
    </style>
    """, unsafe_allow_html=True)


# 👉 Appelle ta fonction avec ton fichier
set_background("wp2932668-noir-wallpaper.jpg")





# ==== Contexte projet ====
ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", ".."))
if ROOT not in sys.path:
    sys.path.insert(0, ROOT)

st.set_page_config(page_title="Mode RC – séance", page_icon="📁", layout="wide")

# Barre de menu haute
c1, c2, c3 = st.columns([1,1,6])
with c1: st.page_link("app.py", label=" Recherche", icon="🔎")
with c2: st.page_link("pages/02_Mode_RC.py", label=" Auto Mém", icon="📁")
with c3: st.write("")

# =========================
#  CONFIG COMMUNE (JSON)
# =========================
CONFIG_DIR = os.path.join(ROOT, "config")
os.makedirs(CONFIG_DIR, exist_ok=True)
APP_CONFIG_FILE = os.path.join(CONFIG_DIR, "app_config.json")




# =========================
#  TF‑IDF local (comme page 1)
# =========================
DEFAULT_INDEX_MEM = os.path.join(ROOT, "chroma_index_mem")

def _get_mem_index_dir() -> str:
    """
    Utilise le même dossier d'index que la page 1 (app.py).
    """
    # Si un autre index vous convient, remplacez par le vôtre :
    return DEFAULT_INDEX_MEM

def _get_collection(persist_dir: str):
    client = chromadb.PersistentClient(path=persist_dir)
    return client.get_or_create_collection(name="memoire_chunks")

@st.cache_resource(show_spinner=False)
def _build_tfidf_index_mem(persist_dir: str):
    """
    Recharge tous les chunks depuis Chroma et construit l’index TF‑IDF local.
    Identique à la logique de la page 1.
    """
    coll = _get_collection(persist_dir)
    batch = coll.get(include=["documents", "metadatas"])
    docs  = (batch.get("documents") or [])
    metas = (batch.get("metadatas") or [])

    docs  = [d if isinstance(d, str)  else ""  for d in docs]
    metas = [m if isinstance(m, dict) else {} for m in metas]

    vect = TfidfVectorizer(
        ngram_range=(1, 2),
        lowercase=True,
        strip_accents="unicode",
        min_df=1,
        max_df=0.9,
    )
    X = vect.fit_transform(docs)
    return {"vect": vect, "X": X, "docs": docs, "metas": metas}

def _tfidf_retrieve_mem(persist_dir: str, query: str, topk: int = 5):
    idx = _build_tfidf_index_mem(persist_dir)
    vect, X, docs, metas = idx["vect"], idx["X"], idx["docs"], idx["metas"]
    if not docs:
        return [], [], []
    qv = vect.transform([query or ""])
    sims = cosine_similarity(qv, X)[0]

    # on prend un “pool” plus large pour dédoublonner ensuite
    top_idx = np.argsort(-sims)[:max(topk * 8, topk)]
    cand = [(int(i), float(sims[i])) for i in top_idx]
    cand.sort(key=lambda t: t[1], reverse=True)

    out_docs, out_metas, out_dists = [], [], []
    for i, sc in cand:
        out_docs.append(docs[i])
        out_metas.append(metas[i])
        out_dists.append(1.0 - sc)  # “distance” simulée
        if len(out_docs) >= topk:
            break
    return out_docs, out_metas, out_dists

def _resolve_source_name(m: dict) -> tuple[str, str, int]:
    """
    Normalise source + nom + page.
    """
    src = (m.get("source") or m.get("source_path") or "").strip()
    name = m.get("source_name") or (os.path.basename(src) if src else "(inconnu)")
    try:
        page = int(m.get("page", 1)); page = page if page >= 1 else 1
    except Exception:
        page = 1
    return src, name, page



# ========= IA helpers (expansion + reranking) =========
def _expand_queries_llm(openai_key: str, brief: str, n: int = 6) -> list[str]:
    """
    Demande au LLM d’élargir la requête : synonymes, normes, procédures.
    Retourne une liste de requêtes courtes (phrases de 3-8 mots).
    """
    try:
        from openai import OpenAI
        cli = OpenAI(api_key=openai_key)
        sys_prompt = (
            "Tu es un assistant d'indexation technique bâtiment/TP.\n"
            "Génère des requêtes courtes et précises (3-8 mots), "
            "incluant synonymes, normes, matériaux, procédés, contrôles, outillage.\n"
            "Réponds en JSON: {\"queries\": [\"...\", \"...\"]}"
        )
        user_prompt = f"Brief:\n{brief}\n\nDonne {n} requêtes différentes et complémentaires."
        resp = cli.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.2,
            response_format={"type":"json_object"},
            messages=[{"role":"system","content":sys_prompt},
                      {"role":"user","content":user_prompt}]
        )
        data = resp.choices[0].message.content or "{}"
        import json as _json
        js = _json.loads(data)
        out = [q.strip() for q in (js.get("queries") or []) if q and q.strip()]
        # fallback si vide
        if not out:
            out = [brief[:200]]
        return out[:max(1, n)]
    except Exception:
        return [brief[:200]]

def _boost_meta_score(meta: dict, base: float, query_tokens: set[str]) -> float:
    """
    Petit boost si le doc a l’air d’un mémoire + si le nom de fichier matche des tokens.
    """
    score = base
    name = (meta.get("source_name") or meta.get("source") or "").lower()
    if "memoire" in name or "mémoire" in name:
        score *= 0.92  # distance plus petite => meilleur rang
    # bonus léger si mot-clé présent dans le nom
    if any(tok in name for tok in query_tokens if tok and len(tok) > 3):
        score *= 0.90
    return score

def _collect_candidates_tfidf(persist_dir: str, queries: list[str], pool: int = 40):
    """
    Lance TF‑IDF pour chaque requête, fusionne, déduplique par (source,page).
    Retourne (docs, metas, dists) triés par distance croissante.
    """
    seen = set()
    pool_docs, pool_metas, pool_dists = [], [], []
    import re as _re
    for q in queries:
        docs, metas, dists = _tfidf_retrieve_mem(persist_dir, q, topk=pool)
        # fusion
        for d, m, dist in zip(docs, metas, dists):
            src = (m.get("source") or m.get("source_path") or "").strip()
            page = m.get("page", 1)
            key = (os.path.normpath(src).lower(), int(page))
            if key in seen:
                continue
            seen.add(key)
            pool_docs.append(d); pool_metas.append(m); pool_dists.append(dist)

    # tri grossier par distance
    items = list(zip(pool_docs, pool_metas, pool_dists))
    items.sort(key=lambda x: x[2])
    return [it[0] for it in items], [it[1] for it in items], [it[2] for it in items]



def _extract_subthemes_llm(openai_key: str, brief: str, max_subthemes: int = 5) -> list[str]:
    """
    Utilise le LLM pour détecter des sous-thèmes distincts dans un chapitre.
    """
    from openai import OpenAI
    cli = OpenAI(api_key=openai_key)
    sys_prompt = (
        "Analyse un brief de cahier des charges technique. "
        "Retourne une liste de sous-thèmes distincts (3-6), concis (1-3 mots chacun). "
        "Ces sous-thèmes doivent couvrir la diversité des points abordés."
    )
    user_prompt = f"Brief:\n{brief}\n\nDonne une liste JSON de sous-thèmes distincts."
    resp = cli.chat.completions.create(
        model="gpt-4o-mini",
        temperature=0.2,
        response_format={"type":"json_object"},
        messages=[{"role":"system","content":sys_prompt},
                  {"role":"user","content":user_prompt}]
    )
    import json
    try:
        data = json.loads(resp.choices[0].message.content)
        subs = [s.strip() for s in (data.get("subthemes") or []) if s.strip()]
        return subs[:max_subthemes] if subs else [brief[:50]]
    except Exception:
        return [brief[:50]]
    
    

def _build_eval_pairs(brief: str, docs: list[str], metas: list[dict], max_chars_doc: int = 900):
    """
    Prépare des paires (brief, extrait) pour le reranking LLM.
    """
    pairs = []
    for d, m in zip(docs, metas):
        clip = (d or "").strip()
        if len(clip) > max_chars_doc:
            # coupe propre sur phrase
            cut = clip[:max_chars_doc]
            last_dot = cut.rfind(".")
            if last_dot > 120:
                cut = cut[:last_dot+1]
            clip = cut + " […]"
        title = m.get("source_name") or os.path.basename(m.get("source") or m.get("source_path") or "") or "(inconnu)"
        page  = m.get("page", 1)
        pairs.append({
            "title": title, "page": int(page),
            "snippet": clip
        })
    return pairs

def _rerank_with_llm(openai_key: str, brief: str, pairs: list[dict], topk: int = 5):
    """
    Demande au LLM de noter chaque pair (0-100) + justification courte.
    Retour: indices des meilleurs candidats dans l’ordre.
    """
    if not pairs:
        return []
    try:
        from openai import OpenAI
        cli = OpenAI(api_key=openai_key)
        sys_prompt = (
            "Tu es un évaluateur de pertinence (bâtiment/TP). "
            "Note chaque extrait (0-100) sur sa pertinence pour le brief. "
            "Réponds JSON: {\"scores\": [{\"i\": idx, \"score\": int, \"why\": \"...\"}, ...]}"
        )
        # On batch si trop long
        batch = pairs[:20]  # sécurité token
        payload = {
            "brief": brief,
            "candidates": [
                {"i": i, "title": p["title"], "page": p["page"], "snippet": p["snippet"]}
                for i, p in enumerate(batch)
            ]
        }
        import json as _json
        user_prompt = _json.dumps(payload, ensure_ascii=False, indent=2)
        resp = cli.chat.completions.create(
            model="gpt-4o-mini",
            temperature=0.0,
            response_format={"type":"json_object"},
            messages=[{"role":"system","content":sys_prompt},
                      {"role":"user","content":user_prompt}]
        )
        data = resp.choices[0].message.content or "{}"
        js = _json.loads(data)
        scores = js.get("scores") or []
        scores.sort(key=lambda x: int(x.get("score", 0)), reverse=True)
        idxs = [int(x.get("i", 0)) for x in scores][:topk]
        return idxs
    except Exception:
        # fallback : sans LLM, on garde l’ordre TF‑IDF
        return list(range(min(topk, len(pairs))))




def _load_app_cfg():
    try:
        if os.path.isfile(APP_CONFIG_FILE):
            with open(APP_CONFIG_FILE, "r", encoding="utf-8") as f:
                return json.load(f) or {}
    except Exception:
        pass
    return {}

def _save_app_cfg(updates: dict):
    cfg = _load_app_cfg()
    cfg.update({k: v for k, v in updates.items() if v is not None})
    with open(APP_CONFIG_FILE, "w", encoding="utf-8") as f:
        json.dump(cfg, f, ensure_ascii=False, indent=2)

def get_openai_key() -> str:
    k = st.session_state.get("openai_key")
    if k is not None:
        return k
    k = (_load_app_cfg().get("openai_key") or "").strip()
    st.session_state.openai_key = k
    return k

def set_openai_key(v: str):
    v = (v or "").strip()
    st.session_state.openai_key = v
    _save_app_cfg({"openai_key": v})

def get_docs_path() -> str:
    p = st.session_state.get("docs_path")
    if p is not None:
        return p
    p = (_load_app_cfg().get("docs_path") or "").strip()
    st.session_state.docs_path = p
    return p

def set_docs_path(v: str):
    v = (v or "").strip()
    st.session_state.docs_path = v
    _save_app_cfg({"docs_path": v})

# =========================
#  Helpers PDF (affichage + résolution)
# =========================
@st.cache_data(show_spinner=False)
def _render_pdf_page_png(abs_path: str, page_index0: int, zoom: float = 2.0) -> bytes:
    doc = fitz.open(abs_path)
    try:
        page_count = doc.page_count
        page_index0 = max(0, min(page_index0, page_count-1))
        page = doc.load_page(page_index0)
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        return pix.tobytes("png")
    finally:
        doc.close()

def _resolve_pdf_path(src: str) -> str | None:
    if not src:
        return None
    try:
        if os.path.isabs(src) and os.path.isfile(src):
            return os.path.normpath(src)
        if os.path.isfile(src):
            return os.path.abspath(src)
        p2 = os.path.normpath(os.path.join(ROOT, src))
        if os.path.isfile(p2):
            return p2
        base = os.path.basename(src)
        d = get_docs_path() or ""
        if d and os.path.isdir(d) and base:
            for r, _, files in os.walk(d):
                if base in files:
                    return os.path.normpath(os.path.join(r, base))
    except Exception:
        pass
    return None

def _read_file_bytes(abs_path: str) -> bytes:
    with open(abs_path, "rb") as f:
        return f.read()

# -----------------------------
# Helpers titres & chapitres
# -----------------------------
_SIMPLE_PREFIX_RX = re.compile(
    r"^(?:\s*\[?\d+\]?\s*|\s*(?:chapitre|section|article)\s*[\d\.]*\s*[:\-\)\u2013\u2014]?\s*)",
    flags=re.IGNORECASE,
)

def _simplify_title(s: str) -> str:
    s = (s or "").strip()
    s = _SIMPLE_PREFIX_RX.sub("", s)
    s = re.sub(r"^[\-–—•·]\s*", "", s)
    s = re.sub(r"\s+", " ", s).strip(" .;:\u2013\u2014-")
    return s

def _explode_mem_chapters_to_points(mem_chaps):
    out = []
    BULLET_RX = re.compile(
        r"(?:^|\n)\s*(?:[-–•·]|[0-9]+\)|\([0-9]+\)|[0-9]+\.)\s*(.+?)\s*(?=\n\s*(?:[-–•·]|[0-9]+\)|\([0-9]+\)|[0-9]+\.)|\Z)",
        flags=re.DOTALL,
    )
    INLINE_AFTER_COLON_RX = re.compile(r":\s*(.+)$", flags=re.DOTALL)
    INLINE_SPLIT_RX = re.compile(r"\s+(?:[-–—•·]|•|·)\s+| ;\s+| – ")

    for c in mem_chaps:
        base_ref   = (c.get("ref") or "").strip()
        base_why   = (c.get("why") or "").strip()
        snippet    = (c.get("snippet") or "").strip()

        points = [m.strip() for m in BULLET_RX.findall(snippet) if m and m.strip()]

        if len(points) <= 1:
            m = INLINE_AFTER_COLON_RX.search(snippet)
            inline_zone = (m.group(1).strip() if m else snippet)
            raw_splits = [p.strip(" .;\n\t\r") for p in INLINE_SPLIT_RX.split(inline_zone) if p.strip()]
            points_inline = [p for p in raw_splits if len(p.split()) >= 3]
            if len(points_inline) >= 2:
                points = points_inline

        if points:
            for p in points:
                simple = _simplify_title(p) or p.strip()
                out.append({"title": simple, "ref": base_ref, "why": base_why, "snippet": p.strip()})
        else:
            base_title = _simplify_title(c.get("title") or "") or _simplify_title(snippet[:120]) or (c.get("title") or "Chapitre")
            out.append({"title": base_title, "ref": base_ref, "why": base_why, "snippet": snippet})
    return out

# --- Outils internes
from app_logic.rc_tools import (
    read_file_text_and_bytes,
    extract_mem_chapters_llm,
    guess_page_for_snippet,
    save_temp_pdf,
    open_pdf_external,
)

# --- Import robuste du backend RC/CCTP
try:
    from app_logic.rc_session import (
        cctp_synthesize_for_chapters_file,  # nouveau (Assistants v2)
        cctp_synthesize_for_chapters_llm,   # ancien (texte)
        search_memoires_for_chapter,
    )
except ImportError:
    from app_logic.rc_session import (
        cctp_synthesize_for_chapters_llm,
        search_memoires_for_chapter,
    )
    cctp_synthesize_for_chapters_file = None  # type: ignore

# -----------------------------
# STATE INIT
# -----------------------------
st.divider()
st.markdown("<h2 style='text-align:center;'>🧭 Process création de mémoire auto / chat gpt</h2>", unsafe_allow_html=True)

# Sidebar — clé OpenAI & options
with st.sidebar:
    st.header("🔐 OpenAI")
    st.session_state.setdefault("openai_key", get_openai_key() or "")
    st.text_input(
        "Clé API OpenAI (sk-…)",
        key="_tmp_openai_key",
        type="password",
        value=(st.session_state.get("openai_key") or ""),
        on_change=lambda: set_openai_key(st.session_state._tmp_openai_key or "")
    )
    if st.session_state.get("openai_key"):
        st.caption("Clé enregistrée ✔️")

    st.divider()
    st.header("Affichage résultats (bouton 4)")
    st.session_state.setdefault("mem_topk", 5)
    st.session_state.mem_topk = st.slider("Nombre de PDF par chapitre (Top-K)", 1, 10, st.session_state.mem_topk)



    st.divider()
    st.header("Nous contacter :")
    "email : neuronalia@gmail.com"
    "tél : 07 67 29 67 94"
    st.caption("en cas de problème contactez nous…")

st.divider()

# Etats principaux
for k, v in [
    ("rc_file_name", ""),
    ("rc_text", ""),
    ("rc_bytes", b""),
    ("rc_saved_path", ""),
    ("cctp_file_name", ""),
    ("cctp_text", ""),
    ("cctp_bytes", b""),
    ("mem_chaps", []),
    ("rc_ch_selected", {}),
    ("cctp_synth_by_chapter", {}),
    ("mem_top1_by_chapter", {}),
    ("pptx_bytes", None),
    ("pptx_filename", ""),
    ("pptx_disk_path", ""),
]:
    if k not in st.session_state:
        st.session_state[k] = v

# -----------------------------
# Uploads (RC / CCTP)
# -----------------------------
col_u1, col_u2 = st.columns(2)
with col_u1:
    rc_upload = st.file_uploader("Déposez le RC/RPAO", type=["pdf","docx","txt","md","rtf"], key="rc_upload")
with col_u2:
    cctp_upload = st.file_uploader("Déposez le CCTP", type=["pdf","docx","txt","md","rtf"], key="cctp_upload")

# Lecture RC
if rc_upload:
    current_name = rc_upload.name
    if current_name != st.session_state.rc_file_name or not st.session_state.rc_bytes:
        data = rc_upload.read()
        rc_text, rc_bytes = read_file_text_and_bytes(current_name, data)
        st.session_state.rc_file_name = current_name
        st.session_state.rc_text = rc_text or ""
        st.session_state.rc_bytes = rc_bytes or b""
        st.session_state.mem_chaps = []
        st.session_state.rc_ch_selected = {}
        st.session_state.rc_saved_path = ""
        st.session_state.cctp_synth_by_chapter = {}
        st.session_state.mem_top1_by_chapter = {}
        st.session_state.pptx_bytes = None
        st.session_state.pptx_filename = ""
        st.session_state.pptx_disk_path = ""

# Lecture CCTP
if cctp_upload:
    current_name = cctp_upload.name
    if current_name != st.session_state.cctp_file_name or not st.session_state.cctp_bytes:
        data = cctp_upload.read()
        cctp_text, cctp_bytes = read_file_text_and_bytes(current_name, data)
        st.session_state.cctp_file_name = current_name
        st.session_state.cctp_text = cctp_text or ""
        st.session_state.cctp_bytes = cctp_bytes or b""
        st.session_state.cctp_synth_by_chapter = {}
        st.session_state.mem_top1_by_chapter = {}
        st.session_state.pptx_bytes = None
        st.session_state.pptx_filename = ""
        st.session_state.pptx_disk_path = ""

# -----------------------------
# Helpers
# -----------------------------
def _require_openai_key_or_stop() -> str:
    key = (st.session_state.get("openai_key") or "").strip()
    if not key:
        st.error("Aucune clé OpenAI détectée. Renseigne-la dans la sidebar (🔐 OpenAI).")
        st.stop()
    return key

# -----------------------------
# Étape 1 — Extraire chapitres RC
# -----------------------------
def _do_extract_rc():
    if not st.session_state.rc_text:
        st.toast("Dépose d’abord le RC/RPAO.", icon="⚠️")
        return
    try:
        key = _require_openai_key_or_stop()
        mem_chaps = extract_mem_chapters_llm(rc_text=st.session_state.rc_text, openai_key=key)
        mem_chaps = _explode_mem_chapters_to_points(mem_chaps)
        st.session_state.mem_chaps = mem_chaps or []
        # IDs stables
        for ch in st.session_state.mem_chaps:
            if "id" not in ch or not ch["id"]:
                ch["id"] = f"{hash((ch.get('title',''), ch.get('snippet','')))}"
        st.session_state.rc_ch_selected = {c["id"]: True for c in st.session_state.mem_chaps}

        if st.session_state.rc_file_name.lower().endswith(".pdf") and st.session_state.rc_bytes:
            st.session_state.rc_saved_path = save_temp_pdf(
                pdf_bytes=st.session_state.rc_bytes,
                root_dir=ROOT,
                suggested_name=None,
            )
        st.toast(f"{len(st.session_state.mem_chaps)} chapitre(s) détecté(s).", icon="✅")
    except Exception as e:
        st.error(f"Erreur d’extraction : {e}")
        st.exception(e)

# -----------------------------
# Étape 2 — RC → CCTP (LLM)  — avec MODE FICHIER préféré
# -----------------------------
def _do_propose_cctp_llm(progress_slot=None):
    if not (st.session_state.cctp_bytes or st.session_state.cctp_text):
        st.toast("Dépose d’abord le CCTP.", icon="⚠️")
        return
    if not st.session_state.mem_chaps:
        st.toast("Extrais d’abord les chapitres du RC.", icon="⚠️")
        return

    key = _require_openai_key_or_stop()

    # chapitres cochés (liste ordonnée)
    selected_titles = []
    for ch in st.session_state.mem_chaps:
        if st.session_state.rc_ch_selected.get(ch["id"], True):
            selected_titles.append(ch["title"])
    if not selected_titles:
        st.toast("Aucun chapitre sélectionné.", icon="⚠️")
        return

    # Barre de progression “par requête”
    slot = progress_slot or st
    ph = slot.empty()
    bar = ph.progress(0, text="Analyse CCTP, cela peut prendre quelques minutes…")

    total = len(selected_titles)
    results = {}
    errors = 0

    prefer_file_mode = bool(cctp_synthesize_for_chapters_file and st.session_state.cctp_bytes)

    for i, title in enumerate(selected_titles, start=1):
        try:
            if prefer_file_mode:
                part = cctp_synthesize_for_chapters_file(
                    rc_chapter_titles=[title],
                    cctp_bytes=st.session_state.cctp_bytes,
                    cctp_name=st.session_state.cctp_file_name or "CCTP.pdf",
                    openai_key=key,
                ) or {}
            else:
                part = cctp_synthesize_for_chapters_llm(
                    rc_chapter_titles=[title],
                    cctp_text=st.session_state.cctp_text,
                    openai_key=key,
                ) or {}
            if isinstance(part, dict):
                results.update(part)
        except Exception as e:
            errors += 1
            # ✅ Affiche la vraie erreur par chapitre
            with st.expander(f"Erreur sur le chapitre « {title} »", expanded=False):
                st.error(str(e))
                st.code("".join(traceback.format_exception(type(e), e, e.__traceback__)))
        finally:
            pct = int(i * 100 / max(total, 1))
            bar.progress(pct, text=f"Analyse CCTP… ({i}/{total})")

    st.session_state.cctp_synth_by_chapter = results
    st.session_state.pptx_bytes = None
    st.session_state.pptx_filename = ""
    st.session_state.pptx_disk_path = ""

    if errors == 0:
        st.toast("Synthèse CCTP par chapitre générée ✔️", icon="🧠")
        bar.progress(100, text="Terminé ✔️")
    else:
        st.toast(f"Synthèse terminée avec {errors} erreur(s).", icon="⚠️")
        bar.progress(100, text=f"Terminé (avec {errors} erreur(s))")

# -----------------------------
# Étape 3 — Export PPTX
# -----------------------------
def _do_export_pptx():
    if not st.session_state.cctp_synth_by_chapter:
        st.toast("Lance d’abord la synthèse CCTP (bouton 2).", icon="⚠️")
        return

    try:
        from pptx import Presentation
        from pptx.util import Pt
        from pptx.enum.text import PP_ALIGN
    except Exception:
        st.error("Le module python-pptx est requis. Installez-le :\n`pip install python-pptx`")
        return

    selected = []
    for ch in st.session_state.mem_chaps:
        if st.session_state.rc_ch_selected.get(ch["id"], True):
            title = ch.get("title", "").strip() or "Chapitre"
            syn = st.session_state.cctp_synth_by_chapter.get(title) or {}
            bullets = syn.get("constraints_synthesis") or []
            selected.append((title, bullets))
    if not selected:
        st.toast("Aucun chapitre sélectionné avec contenu CCTP.", icon="⚠️")
        return

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    s0 = prs.slides.add_slide(title_slide_layout)
    s0.shapes.title.text = "Cadre du mémoire proposé"
    sub = s0.placeholders[1]
    rc_name = st.session_state.get("rc_file_name", "") or "RC"
    cctp_name = st.session_state.get("cctp_file_name", "") or "CCTP"
    sub.text = f"{rc_name}  |  {cctp_name}\nGénéré le {datetime.date.today().strftime('%d/%m/%Y')}"
    for p in sub.text_frame.paragraphs:
        p.font.size = Pt(16)

    content_layout = prs.slide_layouts[1]
    for chap_title, bullets in selected:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = chap_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        if bullets:
            p0 = tf.paragraphs[0]
            p0.text = bullets[0]; p0.level = 0; p0.font.size = Pt(16)
            for b in bullets[1:]:
                p = tf.add_paragraph(); p.text = b; p.level = 0; p.font.size = Pt(16)
        else:
            p = tf.paragraphs[0]; p.text = "Aucune exigence détectée pour ce chapitre."; p.level = 0; p.font.size = Pt(16)

        for p in tf.paragraphs:
            p.alignment = PP_ALIGN.LEFT

    buf = BytesIO(); prs.save(buf); buf.seek(0)
    base = "Cadre_memoires"
    if st.session_state.get("rc_file_name"):
        base = os.path.splitext(st.session_state["rc_file_name"])[0] + "_cadre_memoires"
    fname = f"{base}.pptx"
    st.session_state["pptx_bytes"] = buf.getvalue()
    st.session_state["pptx_filename"] = fname

    try:
        tmp_dir = os.path.join(ROOT, "tmp"); os.makedirs(tmp_dir, exist_ok=True)
        disk_path = os.path.join(tmp_dir, fname)
        with open(disk_path, "wb") as f: f.write(st.session_state["pptx_bytes"])
        st.session_state["pptx_disk_path"] = disk_path
    except Exception:
        st.session_state["pptx_disk_path"] = ""

    st.toast("PowerPoint généré ✔️", icon="📤")

# -----------------------------
# Étape 4 — Recherche mémoires (Top-K)
# -----------------------------
def _do_search_best_memoires():
    if not st.session_state.cctp_synth_by_chapter:
        st.toast("Lance d’abord la synthèse CCTP (bouton 2).", icon="⚠️")
        return

    persist_dir = _get_mem_index_dir()
    if not os.path.isdir(persist_dir):
        st.toast("Index Chroma introuvable : crée d’abord l’index (Mémoires) sur la page 1.", icon="⚠️")
        st.session_state.mem_top_by_chapter = {}
        return

    K = int(st.session_state.get("mem_topk", 5))
    key = _require_openai_key_or_stop()

    results = {}
    chap_processed = chap_found = 0

    # -----------------------------
    #  Fonction interne : extraire sous-thèmes
    # -----------------------------
    def _extract_subthemes_llm(openai_key: str, brief: str, max_subthemes: int = 5) -> list[str]:
        from openai import OpenAI
        cli = OpenAI(api_key=openai_key)
        sys_prompt = (
            "Analyse un brief technique et retourne une liste de sous-thèmes distincts (3-6). "
            "Format JSON strict : {\"subthemes\": [\"...\", \"...\"]}"
        )
        user_prompt = f"Brief:\n{brief}\n\nDonne une liste de sous-thèmes distincts."
        try:
            resp = cli.chat.completions.create(
                model="gpt-4o-mini",
                temperature=0.2,
                response_format={"type":"json_object"},
                messages=[{"role":"system","content":sys_prompt},
                          {"role":"user","content":user_prompt}]
            )
            import json
            data = json.loads(resp.choices[0].message.content)
            subs = [s.strip() for s in (data.get("subthemes") or []) if s.strip()]
            return subs[:max_subthemes] if subs else [brief[:50]]
        except Exception:
            return [brief[:50]]

    # -----------------------------
    #  Parcours des chapitres
    # -----------------------------
    for ch in st.session_state.mem_chaps:
        if not st.session_state.rc_ch_selected.get(ch["id"], True):
            continue

        title = (ch.get("title") or "").strip()
        syn = st.session_state.cctp_synth_by_chapter.get(title) or {}
        bullets = (syn.get("constraints_synthesis") or [])[:12]
        why_rc = (ch.get("why") or "").strip()

        brief_parts = [f"Chapitre: {title}"]
        if why_rc: brief_parts.append(f"Contexte RC: {why_rc}")
        if bullets: brief_parts.append("Exigences CCTP: " + " ; ".join(bullets))
        brief = "\n".join(brief_parts)

        # ---- 1) Génération des sous-thèmes ----
        subthemes = _extract_subthemes_llm(key, brief, max_subthemes=5)

        # ---- 2) Recherche par sous-thème ----
        candidates = []
        for stheme in subthemes:
            try:
                docs, metas, dists = _tfidf_retrieve_mem(persist_dir, stheme, topk=K*4)
                for d, m, dist in zip(docs, metas, dists):
                    candidates.append((d, m, dist, stheme))
            except Exception as e:
                st.error(f"Erreur recherche TF-IDF pour '{stheme}' : {e}")

        # ---- 3) Fusion + diversité ----
        # ---- 3) Fusion + diversité (pool K élargi) ----
        # On crée un pool diversifié (max 2 par sous-thème), puis on laissera le LLM reranker ce pool.
        seen = {}
        pool_div = []
        POOL_TARGET = max(K * 3, K + 3)
        for d, m, dist, stheme in sorted(candidates, key=lambda x: x[2]):
            src = (m.get("source") or "").strip()
            page = m.get("page", 1)
            key_src = (os.path.normpath(src).lower(), int(page))
            if key_src in seen:
                continue
            if sum(1 for _,_,_,s in pool_div if s == stheme) >= 2:
                continue
            pool_div.append((d, m, dist, stheme))
            seen[key_src] = True
            if len(pool_div) >= POOL_TARGET:
                break

        # ---- 4) RERANK LLM (pertinence au brief) ----
        if pool_div:
            # On rerank le pool via LLM
            pairs = _build_eval_pairs(
                brief=brief,
                docs=[x[0] for x in pool_div],
                metas=[x[1] for x in pool_div],
                max_chars_doc=900
            )
            order = _rerank_with_llm(
                openai_key=key,
                brief=brief,
                pairs=pairs,
                topk=K
            )
            # fallback si le LLM ne répond pas
            if not order:
                order = list(range(min(K, len(pool_div))))

            # Reconstitution des items selon l'ordre LLM, en gardant unicité (source,page)
            picked, seen2 = [], set()
            for i in order:
                if i < 0 or i >= len(pool_div): 
                    continue
                d, m, dist, stheme = pool_div[i]
                src = (m.get("source") or "").strip()
                page = int(m.get("page", 1))
                key_src = (os.path.normpath(src).lower(), page)
                if key_src in seen2:
                    continue
                picked.append((d, m, dist, stheme))
                seen2.add(key_src)
                if len(picked) >= K:
                    break
        else:
            picked = []

        # ---- 5) Stockage ----
        if picked:
            chap_found += 1
            items = []
            for d, m, dist, stheme in picked:
                src, name, page = _resolve_source_name(m)
                items.append({
                    "doc": d,
                    "meta": {"source": src, "source_name": name, "page": page, "subtheme": stheme},
                    "dist": dist
                })
            results[title] = items
        else:
            results[title] = []


        chap_processed += 1

    st.session_state.mem_top_by_chapter = results

    # ---- Feedback utilisateur ----
    if chap_processed == 0:
        st.toast("Aucun chapitre sélectionné.", icon="⚠️")
    elif chap_found == 0:
        st.toast("Aucun résultat pertinent trouvé. Vérifie l’indexation.", icon="⚠️")
    else:
        st.toast(f"Recherche IA diversifiée terminée ✔️ ({chap_found}/{chap_processed} chapitres)", icon="🔎")

# -----------------------------
# UI — Ordre des actions
# -----------------------------
if st.button("1) 1- Extraire les chapitres du RC/RPAO en lien avec le mémoire technique", key="btn_extract_rc"):
    with st.spinner("Extraction des chapitres en cours..."):
        _do_extract_rc()

st.divider()
st.markdown("<h3 style='text-align:center;'>Chapitres détectés</h3>", unsafe_allow_html=True)

if st.session_state.mem_chaps:
    ids_a_supprimer = set()
    for pos, chap in enumerate(st.session_state.mem_chaps, start=1):
        chap_id = chap["id"]
        clean_title = re.sub(r"^\[?\d+\]?\s*", "", chap.get('title', '')).strip()
        clean_title = re.sub(r"^\d+(?:\.\d+)*\s*", "", clean_title)

        with st.container(border=True):
            col1, col2, col3 = st.columns([8, 2, 2])
            with col1:
                st.markdown(f"**{clean_title}**")
            with col2:
                checked = st.checkbox(
                    "Sélectionner",
                    key=f"rc_ch_sel_{chap_id}",
                    value=st.session_state.rc_ch_selected.get(chap_id, True),
                    help="Inclure ce chapitre dans la suite",
                )
                st.session_state.rc_ch_selected[chap_id] = checked
                if not checked:
                    ids_a_supprimer.add(chap_id)
            with col3:
                if st.session_state.rc_file_name.lower().endswith(".pdf"):
                    def _open_at(idx=pos):
                        mem = st.session_state.mem_chaps[idx-1]
                        title = mem.get("title","")
                        snippet = mem.get("snippet","")
                        page_guess = 1
                        if st.session_state.rc_bytes and st.session_state.rc_file_name.lower().endswith(".pdf"):
                            page_guess = guess_page_for_snippet(
                                st.session_state.rc_bytes, title, snippet, mem.get("ref","")
                            )
                        path = st.session_state.rc_saved_path
                        if not path or not os.path.isfile(path):
                            path = save_temp_pdf(st.session_state.rc_bytes, ROOT, None)
                            st.session_state.rc_saved_path = path
                        open_pdf_external(path, page_guess, st.session_state.get("custom_pdf_viewer", ""))
                    st.button("🔗 Ouvrir le PDF", key=f"open_rc_{chap_id}", on_click=_open_at)
                else:
                    st.caption("(PDF RC non chargé)")
    if ids_a_supprimer:
        st.session_state.mem_chaps = [c for c in st.session_state.mem_chaps if c["id"] not in ids_a_supprimer]
        st.rerun()
else:
    st.caption("Les chapitres apparaîtront ici après extraction ou ajout manuel.")

with st.form(key="manual_add_form", clear_on_submit=True):
    man_title = st.text_input(
        "Ajouter un chapitre manuellement",
        placeholder="Équipements et matériels utilisés, moyens de translation horizontal & vertical"
    )
    submitted = st.form_submit_button("➕ Ajouter à la liste")
    if submitted:
        simple_title = _simplify_title(man_title)
        if not simple_title:
            st.warning("Le titre est obligatoire.")
        else:
            chap_id = str(uuid.uuid4())
            st.session_state.mem_chaps.append({
                "id": chap_id,
                "title": simple_title,
                "ref": "",
                "why": "",
                "snippet": man_title.strip(),
            })
            st.session_state.rc_ch_selected[chap_id] = True
            st.toast(f"Chapitre ajouté : {simple_title}", icon="✅")
            st.rerun()

# 2) Bouton 2 — Synthèse CCTP
cctp_btn_col = st.container()
cctp_progress_row = st.container()
with cctp_btn_col:
    clicked = st.button(
        "2) 2- Rechercher depuis le CCTP les articles en lien avec le mémoire technique",
        key="btn_cctp_llm"
    )
if clicked:
    _do_propose_cctp_llm(progress_slot=cctp_progress_row)

st.divider()
if st.session_state.cctp_synth_by_chapter:
    st.subheader("Exigences CCTP par chapitre (synthèse)")
    for ch in st.session_state.mem_chaps:
        if not st.session_state.rc_ch_selected.get(ch["id"], True):
            continue
        title = ch["title"]
        syn = st.session_state.cctp_synth_by_chapter.get(title)
        if not syn:
            continue
        bullets = syn.get("constraints_synthesis") or []
        cov = syn.get("coverage")
        with st.expander(f"📌 {title}", expanded=False):
            if bullets:
                for b in bullets: st.write(b)
            else:
                st.caption("Aucune exigence détectée pour ce chapitre.")
            if isinstance(cov, (int, float)):
                st.caption(f"Couverture estimée : {cov:.0%}")

# 3) Bouton 3 — Export PPTX
st.button("3) 3- Exporter en PWP le cadre du mémoire proposé", key="btn_PWP", on_click=_do_export_pptx)
if st.session_state.get("pptx_bytes"):
    col_dl1, _ = st.columns([3, 2])
    with col_dl1:
        st.download_button(
            "⬇️ Télécharger le PowerPoint",
            data=st.session_state["pptx_bytes"],
            file_name=st.session_state.get("pptx_filename", "Cadre_memoires.pptx"),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="dl_pptx_main"
        )

st.divider()

# 4) Bouton 4 — Recherche mémoires
st.button(
    "4) 4- Rechercher le mémoire technique tmh correspondant au chapitre séléctionné",
    key="btn_find_best",
    on_click=_do_search_best_memoires
)

# Résultats Top-K
if st.session_state.get("mem_top_by_chapter"):
    st.subheader("Mémoires correspondants par chapitre (Top-K)")
    K = int(st.session_state.get("mem_topk", 5))
    for ch in st.session_state.mem_chaps:
        if not st.session_state.rc_ch_selected.get(ch["id"], True):
            continue
        title = ch["title"]
        items = st.session_state.mem_top_by_chapter.get(title) or []
        with st.expander(f"📄 {title} — {len(items)} PDF affiché(s)", expanded=False):
            if not items:
                st.caption("Aucun résultat pour ce chapitre.")
                continue
            N_PER_ROW = 2
            for i in range(0, len(items), N_PER_ROW):
                cols = st.columns(N_PER_ROW)
                for j, col in enumerate(cols):
                    idx = i + j
                    if idx >= len(items): 
                        with col: st.empty(); continue
                    entry = items[idx]
                    meta  = entry.get("meta") or {}
                    src   = meta.get("source_path") or meta.get("source") or ""
                    name  = meta.get("source_name") or os.path.basename(src) or "(inconnu)"
                    page  = meta.get("page", 1)
                    doc_txt = entry.get("doc") or ""
                    abs_path = _resolve_pdf_path(src)
                    with col:
                        st.markdown(f"**[{idx+1}] {name} — page {page}**")
                        if not abs_path or not os.path.isfile(abs_path):
                            st.warning("Fichier introuvable (chemin index obsolète et non relocalisable).")
                        else:
                            try:
                                png_bytes = _render_pdf_page_png(abs_path, int(page)-1, zoom=2.0)
                                st.image(png_bytes, use_container_width=True)
                            except Exception as e:
                                st.error(f"Impossible d’afficher la page : {e}")
                            if st.button("🔗 Ouvrir le PDF", key=f"open_mem_{title}_{idx}"):
                                try:
                                    pdf_bytes = _read_file_bytes(abs_path)
                                    page_guess = guess_page_for_snippet(
                                        pdf_bytes=pdf_bytes,
                                        title=title,
                                        snippet=doc_txt,
                                        ref=meta.get("ref","")
                                    ) or page
                                except Exception:
                                    page_guess = page
                                open_pdf_external(abs_path, page_guess, st.session_state.get("custom_pdf_viewer",""))
